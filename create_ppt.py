import glob, sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0xD4, 0xFF)
ACCENT_PURPLE = RGBColor(0xA8, 0x55, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_GRAY = RGBColor(0x33, 0x33, 0x44)
GOLD = RGBColor(0xFF, 0xD7, 0x00)
GREEN = RGBColor(0x00, 0xE6, 0x76)
RED = RGBColor(0xFF, 0x45, 0x45)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)

def add_bg(slide, color=DARK_BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf

def add_para(tf, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return p

def add_accent_bar(slide, left, top, width=0.3, height=3):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(left), Cm(top), Cm(width), Cm(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()

def add_section_title(slide, title, subtitle=""):
    add_bg(slide)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Cm(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()
    add_text(slide, 2, 1, 30, 2, title, size=32, color=ACCENT_BLUE, bold=True)
    if subtitle:
        add_text(slide, 2, 3.5, 30, 1.5, subtitle, size=16, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 1: 표지
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

accent_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Cm(1), prs.slide_height)
accent_rect.fill.solid()
accent_rect.fill.fore_color.rgb = ACCENT_BLUE
accent_rect.line.fill.background()

add_text(slide, 2, 1.5, 10, 1.5, "K-Digital Training", size=16, color=ACCENT_BLUE, bold=True)
add_text(slide, 2, 4, 28, 3, "INFINITY LEVEL UP", size=54, color=WHITE, bold=True)

tf = add_text(slide, 2, 8, 28, 2, "초간편 하이퍼 캐주얼 2D RPG", size=24, color=LIGHT_GRAY)
add_para(tf, "팀별 프로젝트 결과보고서", size=20, color=LIGHT_GRAY)

add_text(slide, 2, 13, 28, 1, "장르: 하이퍼 캐주얼 RPG  |  타겟: 10세 이상 전연령  |  플랫폼: 모바일", size=14, color=ACCENT_BLUE)
add_text(slide, 2, 14.5, 28, 1.5, "Unity  |  C#  |  Backend SDK  |  2D RPG", size=16, color=LIGHT_GRAY)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Cm(18.5), prs.slide_width, Cm(0.6))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_BLUE
bar.line.fill.background()

# ============================================================
# 슬라이드 2: 목차
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "목차", "TABLE OF CONTENTS")

items_list = [
    ("01", "프로젝트 개요"),
    ("02", "팀 구성 및 역할 분담"),
    ("03", "프로젝트 수행 절차"),
    ("04", "프로젝트 소개 및 핵심 기능"),
    ("05", "모바일 게임 시장 분석"),
    ("06", "타 게임 비교 분석 & 차별점"),
    ("07", "기술 스택"),
    ("08", "핵심 기능 구현 내용"),
    ("09", "치매예방 미니게임"),
    ("10", "아키텍처 & 설계 패턴"),
    ("11", "QA 결과 및 개선 사항"),
    ("12", "시연 동영상"),
]

for idx, (num, title) in enumerate(items_list):
    y = 5.0 + idx * 1.1
    add_text(slide, 3, y, 3, 1.1, num, size=24, color=ACCENT_BLUE, bold=True)
    add_text(slide, 6, y, 20, 1.1, title, size=17, color=WHITE)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(3), Cm(y+1.05), Cm(25), Cm(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x44, 0x44, 0x66)
    line.line.fill.background()

# ============================================================
# 슬라이드 3: 프로젝트 개요
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "01  프로젝트 개요")

add_accent_bar(slide, 2, 5, 0.2, 12)

tf = add_text(slide, 3, 5, 14, 1, "프로젝트 명", size=14, color=ACCENT_BLUE, bold=True)
add_para(tf, "INFINITY LEVEL UP", size=24, color=WHITE, bold=True)
add_para(tf, "", size=8)
add_para(tf, "장르", size=14, color=ACCENT_BLUE, bold=True)
add_para(tf, "초간편 하이퍼 캐주얼 2D RPG", size=18, color=WHITE)
add_para(tf, "", size=8)
add_para(tf, "타겟층", size=14, color=ACCENT_BLUE, bold=True)
add_para(tf, "10세 이상 전연령", size=18, color=WHITE)
add_para(tf, "", size=8)
add_para(tf, "타겟 플랫폼", size=14, color=ACCENT_BLUE, bold=True)
add_para(tf, "모바일 (Android / iOS)", size=18, color=WHITE)
add_para(tf, "", size=8)
add_para(tf, "레퍼런스 게임", size=14, color=ACCENT_BLUE, bold=True)
add_para(tf, "The Tower  |  버섯커키우기  |  레알팜", size=18, color=WHITE)
add_para(tf, "", size=8)
add_para(tf, "기획 의도", size=14, color=ACCENT_BLUE, bold=True)
add_para(tf, "누구나 쉽게 즐길 수 있는 초간편 조작으로", size=16, color=WHITE)
add_para(tf, "방치형 RPG의 성장 재미를 극대화", size=16, color=WHITE)

# Right - key features
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(19), Cm(5), Cm(13), Cm(12))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3E)
box.line.color.rgb = RGBColor(0x44, 0x44, 0x66)

tf = add_text(slide, 20, 5.5, 11, 1, "핵심 기능", size=18, color=GOLD, bold=True)
features = [
    "실시간 2D 전투 (근거리/원거리)",
    "무한의탑 (끝없는 도전 콘텐츠)",
    "월드보스 레이드 (협동 전투)",
    "장비 & 동료 가챠 시스템",
    "장비 강화/레벨업 시스템",
    "농장 미니게임 (재배, 건물, 경매)",
    "스킬 트리 & 핫바",
    "치매예방 미니게임 (두뇌 훈련)",
    "VIP / 랭킹 / 메일 / 채팅",
    "Backend 서버 연동",
]
for f in features:
    add_para(tf, f"  {f}", size=14, color=WHITE)

# ============================================================
# 슬라이드 4: 팀 구성
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "02  팀 구성 및 역할 분담")

table_shape = slide.shapes.add_table(5, 3, Cm(2), Cm(5.5), Cm(30), Cm(10))
table = table_shape.table
table.columns[0].width = Cm(6)
table.columns[1].width = Cm(6)
table.columns[2].width = Cm(18)

headers = ["훈련생", "역할", "담당 업무"]
rows_data = [
    ["OOO", "팀장", "전투 시스템, 게임 매니저, 씬 전환, 저장/로드"],
    ["OOO", "팀원", "가챠, 인벤토리, 장비 강화, 제작 시스템"],
    ["OOO", "팀원", "농장, 퀘스트, 업적, UI/UX"],
    ["OOO", "멘토", "코드 리뷰, 아키텍처 설계 자문"],
]

for ci, h in enumerate(headers):
    cell = table.cell(0, ci)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0x00, 0x7A, 0xCC)

for ri, row_data in enumerate(rows_data):
    for ci, val in enumerate(row_data):
        cell = table.cell(ri+1, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER if ci < 2 else PP_ALIGN.LEFT
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x44) if ri % 2 == 0 else RGBColor(0x22, 0x22, 0x3E)

# ============================================================
# 슬라이드 5: 프로젝트 수행 절차
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "03  프로젝트 수행 절차")

add_text(slide, 2, 4, 30, 1, "프로젝트 개발 일정", size=14, color=LIGHT_GRAY)

table_shape = slide.shapes.add_table(7, 4, Cm(2), Cm(6), Cm(30), Cm(11))
table = table_shape.table
table.columns[0].width = Cm(6)
table.columns[1].width = Cm(7)
table.columns[2].width = Cm(11)
table.columns[3].width = Cm(6)

schedule = [
    ["구분", "기간", "활동", "비고"],
    ["사전 기획", "1주차", "게임 컨셉, 레퍼런스 분석, 기획서", "아이디어 선정"],
    ["프로토타입", "2~3주차", "전투, 캐릭터 컨트롤러, 기본 UI", "핵심 메카닉"],
    ["콘텐츠 개발", "4~5주차", "가챠, 인벤토리, 강화, 스킬, 농장", "주요 시스템"],
    ["서버 연동", "5~6주차", "Backend SDK (인증, 채팅, 랭킹, 저장)", "백엔드 통합"],
    ["QA/수정", "6~7주차", "버그 수정, 밸런싱, UI/UX 개선", "품질 관리"],
    ["총 개발기간", "약 7주", "기획 -> 개발 -> QA -> 완성", ""],
]

for ri, row_data in enumerate(schedule):
    for ci, val in enumerate(row_data):
        cell = table.cell(ri, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE
            p.font.bold = (ri == 0)
            p.alignment = PP_ALIGN.CENTER
        if ri == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x00, 0x7A, 0xCC)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x44) if ri % 2 == 0 else RGBColor(0x22, 0x22, 0x3E)

# ============================================================
# 슬라이드 6: 프로젝트 소개
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "04  프로젝트 소개")

add_accent_bar(slide, 2, 5, 0.2, 13)

tf = add_text(slide, 3, 5, 28, 14, "INFINITY LEVEL UP", size=26, color=WHITE, bold=True)
add_para(tf, "", size=6)
add_para(tf, "초간편 하이퍼 캐주얼 2D RPG  |  10세 이상 전연령  |  모바일", size=16, color=ACCENT_BLUE)
add_para(tf, "", size=6)
add_para(tf, "레퍼런스: The Tower(방치형 전투) + 버섯커키우기(성장/수집) + 레알팜(농장 경영)", size=14, color=GOLD)
add_para(tf, "", size=8)
add_para(tf, "게임 소개", size=18, color=GOLD, bold=True)
add_para(tf, "INFINITY LEVEL UP은 누구나 손쉽게 즐길 수 있는 초간편 조작의", size=15, color=WHITE)
add_para(tf, "하이퍼 캐주얼 RPG입니다. The Tower처럼 간단한 방치형 전투,", size=15, color=WHITE)
add_para(tf, "버섯커키우기처럼 다양한 수집/성장 요소, 레알팜처럼 힐링 농장을", size=15, color=WHITE)
add_para(tf, "결합하여 끝없는 성장을 경험할 수 있습니다.", size=15, color=WHITE)
add_para(tf, "", size=6)
add_para(tf, "주요 콘텐츠", size=18, color=GOLD, bold=True)
add_para(tf, "  - 원터치 자동 전투 (근거리/원거리 캐릭터 선택)", size=14, color=WHITE)
add_para(tf, "  - 무한의탑: 층수 도전형 엔드 콘텐츠", size=14, color=WHITE)
add_para(tf, "  - 월드보스 레이드: 대형 보스 협동 전투", size=14, color=WHITE)
add_para(tf, "  - 장비 & 동료 가챠 (5티어 확률 시스템)", size=14, color=WHITE)
add_para(tf, "  - 장비 강화 / 농장 / 치매예방 미니게임", size=14, color=WHITE)
add_para(tf, "  - 스킬 트리, 채팅, 랭킹, VIP, 오프라인 보상", size=14, color=WHITE)

# ============================================================
# 슬라이드 7: 모바일 게임 시장 분석 (NEW!)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "05  모바일 게임 시장 분석", "Market Analysis")

# Left: Global market
add_accent_bar(slide, 2, 5, 0.2, 5.5)
tf = add_text(slide, 3, 5, 14, 5.5, "글로벌 모바일 게임 시장", size=20, color=GOLD, bold=True)
add_para(tf, "", size=4)
add_para(tf, "  2024년 시장 규모: 약 $920억 (USD)", size=15, color=WHITE)
add_para(tf, "  전체 게임 시장의 약 49% 차지", size=15, color=WHITE)
add_para(tf, "  연평균 성장률(CAGR): 7~8%", size=15, color=WHITE)
add_para(tf, "  2028년 예상: $1,200억 이상", size=15, color=WHITE)
add_para(tf, "  (출처: Newzoo, Sensor Tower)", size=11, color=LIGHT_GRAY)

# Left: Korean market
add_accent_bar(slide, 2, 11.5, 0.2, 5.5)
tf = add_text(slide, 3, 11.5, 14, 5.5, "한국 모바일 게임 시장", size=20, color=GOLD, bold=True)
add_para(tf, "", size=4)
add_para(tf, "  2024년 시장 규모: 약 $65~70억 (USD)", size=15, color=WHITE)
add_para(tf, "  세계 4위 모바일 게임 시장", size=15, color=WHITE)
add_para(tf, "  Google Play 점유율: 약 70~75%", size=15, color=WHITE)
add_para(tf, "  캐주얼/방치형 RPG 강세", size=15, color=WHITE)
add_para(tf, "  (출처: KOCCA, data.ai)", size=11, color=LIGHT_GRAY)

# Right: Market size bar chart (text-based)
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(19), Cm(5), Cm(13), Cm(12.5))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3E)
box.line.color.rgb = RGBColor(0x44, 0x44, 0x66)

tf = add_text(slide, 20, 5.5, 11, 12, "글로벌 모바일 게임 시장 성장 추이", size=15, color=ACCENT_BLUE, bold=True)
add_para(tf, "", size=6)

# Bar chart using shapes
market_data = [
    ("2022", 77, RGBColor(0x33, 0x77, 0xBB)),
    ("2023", 83, RGBColor(0x44, 0x88, 0xCC)),
    ("2024", 92, RGBColor(0x55, 0x99, 0xDD)),
    ("2025(E)", 99, RGBColor(0x00, 0xAA, 0xEE)),
    ("2026(E)", 107, RGBColor(0x00, 0xBB, 0xFF)),
    ("2028(E)", 120, ACCENT_BLUE),
]

for idx, (year, value, color) in enumerate(market_data):
    y_pos = 7.5 + idx * 1.5
    bar_width = value / 120 * 10  # Scale to max 10cm

    # Year label
    add_text(slide, 20, y_pos, 3, 1, year, size=11, color=WHITE, bold=True)

    # Bar
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(23.5), Cm(y_pos+0.1), Cm(bar_width), Cm(0.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    # Value label
    add_text(slide, 23.5 + bar_width + 0.2, y_pos, 3, 1, f"${value}B", size=10, color=WHITE)

add_text(slide, 20, 17, 11, 1, "(단위: 10억 USD, E=예상)", size=10, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 8: 하이퍼 캐주얼 & 방치형 시장 (NEW!)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "05  하이퍼 캐주얼 & 방치형 RPG 시장", "장르별 시장 가능성")

# Hyper casual
add_accent_bar(slide, 2, 5, 0.2, 5)
tf = add_text(slide, 3, 5, 14, 5, "하이퍼 캐주얼 시장 트렌드", size=20, color=GOLD, bold=True)
add_para(tf, "", size=4)
add_para(tf, "  하이퍼 캐주얼 -> '하이브리드 캐주얼' 진화 중", size=14, color=WHITE)
add_para(tf, "  단순 조작 + 깊은 메타(성장/수집) 결합이 대세", size=14, color=WHITE)
add_para(tf, "  CPI(설치당 비용) 최저 수준 ($0.2~0.5)", size=14, color=WHITE)
add_para(tf, "  전연령 접근성으로 유저 확보 용이", size=14, color=WHITE)

# Idle RPG
add_accent_bar(slide, 2, 11, 0.2, 5)
tf = add_text(slide, 3, 11, 14, 5, "방치형/Idle RPG 시장", size=20, color=GOLD, bold=True)
add_para(tf, "", size=4)
add_para(tf, "  아시아(한국/일본/중국) 시장에서 강세", size=14, color=WHITE)
add_para(tf, "  한국 Google Play 매출 상위권 상주", size=14, color=WHITE)
add_para(tf, "  방치 + 가챠 결합 모델 성장세", size=14, color=WHITE)
add_para(tf, "  바쁜 현대인에게 최적화된 플레이 패턴", size=14, color=WHITE)

# Right: why our game fits
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(19), Cm(5), Cm(13), Cm(12))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3E)
box.line.color.rgb = RGBColor(0x44, 0x44, 0x66)

tf = add_text(slide, 20, 5.5, 11, 11, "INFINITY LEVEL UP의 시장 적합성", size=16, color=ACCENT_BLUE, bold=True)
add_para(tf, "", size=6)

fits = [
    ("하이브리드 캐주얼 트렌드 부합", "초간편 조작 + 깊은 성장/수집 메타\n하이퍼 캐주얼의 진화 방향과 일치"),
    ("방치형 RPG 시장 성장", "The Tower/버섯커키우기와 동일 장르\n검증된 시장에서의 차별화 경쟁"),
    ("전연령 접근성", "10세 이상 타겟으로 넓은 유저풀\n가족 단위 플레이 가능"),
    ("다중 수익 모델", "가챠 + VIP + 광고 수익\n다양한 수익원으로 안정적 BM"),
]

for title, desc in fits:
    add_para(tf, f"  {title}", size=14, color=GREEN, bold=True)
    for line in desc.split("\n"):
        add_para(tf, f"    {line}", size=12, color=WHITE)
    add_para(tf, "", size=4)

# ============================================================
# 슬라이드 9: 타 게임 비교 분석 (The Tower, 버섯커키우기, 레알팜)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "06  타 게임 비교 분석", "레퍼런스 게임과의 비교")

table_shape = slide.shapes.add_table(13, 5, Cm(1.5), Cm(4.5), Cm(31), Cm(14))
table = table_shape.table
table.columns[0].width = Cm(6)
table.columns[1].width = Cm(7)
table.columns[2].width = Cm(6)
table.columns[3].width = Cm(6)
table.columns[4].width = Cm(6)

compare_data = [
    ["비교 항목", "INFINITY\nLEVEL UP", "The Tower", "버섯커\n키우기", "레알팜"],
    ["장르", "하이퍼 캐주얼\nRPG", "방치형\n타워 디펜스", "방치형\nRPG", "농장\n시뮬레이션"],
    ["타겟층", "10세+ 전연령", "전연령", "전연령", "전연령"],
    ["조작 난이도", "초간편\n(원터치)", "초간편\n(원터치)", "초간편\n(원터치)", "간편\n(터치)"],
    ["전투 시스템", "O (실시간\n2D 전투)", "O (자동\n타워 공격)", "O (자동\n전투)", "X"],
    ["수집/가챠", "O (장비+동료\n5티어 가챠)", "O (무기\n업그레이드)", "O (캐릭터\n가챠)", "X"],
    ["농장/건설", "O (농장+건물\n+경매장)", "X", "X", "O (농장\n경영)"],
    ["장비 강화", "O (30레벨\n성공률)", "O (강화)", "O (강화)", "X"],
    ["소셜 기능", "O (채팅/랭킹\n/메일)", "O (랭킹)", "O (길드)", "O (마을)"],
    ["무한의탑", "O (층수 도전\n엔드 콘텐츠)", "O (층수\n도전)", "X", "X"],
    ["월드보스", "O (협동\n레이드)", "X", "X", "X"],
    ["미니게임", "O (치매예방\n두뇌 훈련)", "X", "X", "X"],
]

col_colors = [
    RGBColor(0x33, 0x33, 0x55),
    RGBColor(0x00, 0x7A, 0xCC),
    RGBColor(0x55, 0x44, 0x22),
    RGBColor(0x44, 0x22, 0x55),
    RGBColor(0x22, 0x55, 0x33),
]

for ri, row_data in enumerate(compare_data):
    for ci, val in enumerate(row_data):
        cell = table.cell(ri, ci)
        cell.text = val
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
            p.font.bold = (ri == 0 or ci == 0)
            p.alignment = PP_ALIGN.CENTER
        if ri == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = col_colors[ci]
        elif ci == 1:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x00, 0x3D, 0x66)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x44) if ri % 2 == 0 else RGBColor(0x22, 0x22, 0x3E)

# ============================================================
# 슬라이드 10: 게임성 비교 점수 차트
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "06  게임성 비교", "항목별 점수 비교 (5점 만점)")

games = ["INFINITY LEVEL UP", "The Tower", "버섯커키우기", "레알팜"]
game_colors = [ACCENT_BLUE, ORANGE, ACCENT_PURPLE, GREEN]

# Legend
for gi, (game, gcolor) in enumerate(zip(games, game_colors)):
    x = 3 + gi * 8
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x), Cm(4.5), Cm(0.8), Cm(0.8))
    box.fill.solid()
    box.fill.fore_color.rgb = gcolor
    box.line.fill.background()
    add_text(slide, x+1.2, 4.4, 6, 1, game, size=12, color=WHITE, bold=True)

categories = [
    ("조작 편의성",   5, 5, 5, 4),
    ("수집 다양성",   5, 3, 4, 2),
    ("전투 재미",     4, 4, 3, 1),
    ("성장 시스템",   5, 4, 5, 2),
    ("부가 콘텐츠",   5, 2, 3, 4),
    ("소셜 기능",     5, 2, 3, 3),
    ("접근성(연령)", 5, 5, 5, 5),
]

table_shape = slide.shapes.add_table(8, 5, Cm(2), Cm(6), Cm(30), Cm(11))
table = table_shape.table
table.columns[0].width = Cm(6)
table.columns[1].width = Cm(6)
table.columns[2].width = Cm(6)
table.columns[3].width = Cm(6)
table.columns[4].width = Cm(6)

header = ["항목", "INFINITY\nLEVEL UP", "The Tower", "버섯커\n키우기", "레알팜"]
for ci, h in enumerate(header):
    cell = table.cell(0, ci)
    cell.text = h
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    if ci == 1:
        cell.fill.fore_color.rgb = RGBColor(0x00, 0x7A, 0xCC)
    else:
        cell.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x55)

for ri, (cat, s1, s2, s3, s4) in enumerate(categories):
    scores = [s1, s2, s3, s4]
    cell = table.cell(ri+1, 0)
    cell.text = cat
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x44)

    for ci, score in enumerate(scores):
        cell = table.cell(ri+1, ci+1)
        stars = "\u2605" * score + "\u2606" * (5 - score)
        cell.text = stars
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.color.rgb = game_colors[ci]
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        if ci == 0:
            cell.fill.fore_color.rgb = RGBColor(0x00, 0x3D, 0x66)
        else:
            cell.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x44) if ri % 2 == 0 else RGBColor(0x22, 0x22, 0x3E)

# ============================================================
# 슬라이드 11: 차별화 포인트
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "06  차별화 포인트", "INFINITY LEVEL UP만의 강점")

cards = [
    ("The Tower에서 가져온 것", [
        "  초간편 방치형 전투 메카닉",
        "  원터치 조작의 접근성",
        "  자동 전투 + 오프라인 보상",
    ], "차별화: 2D RPG 전투 + 스킬 시스템 추가", ACCENT_BLUE),
    ("버섯커키우기에서 가져온 것", [
        "  다양한 수집/성장 요소",
        "  가챠 기반 캐릭터/장비 수집",
        "  장비 강화 시스템",
    ], "차별화: 5티어 가챠 + 30레벨 강화 + 동료 시스템", GOLD),
    ("레알팜에서 가져온 것", [
        "  농장 경영 시스템",
        "  작물 재배/수확 사이클",
        "  힐링 콘텐츠로 리텐션 강화",
    ], "차별화: 경매장(PvP 거래) + 건물 건설 추가", GREEN),
    ("INFINITY만의 독자 요소", [
        "  치매예방 미니게임 (두뇌 훈련)",
        "  실시간 5채널 채팅 시스템",
        "  VIP + 랭킹 + 스킬 트리",
        "  Backend 클라우드 세이브",
    ], "3개 게임 장점 통합 + 치매예방 + 소셜", ACCENT_PURPLE),
]

for idx, (title, items, diff, color) in enumerate(cards):
    col = idx % 2
    row = idx // 2
    left = 2 + col * 16.5
    top = 5 + row * 6.5

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(left), Cm(top), Cm(14.5), Cm(5.8))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3E)
    box.line.color.rgb = RGBColor(0x44, 0x44, 0x66)

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(left), Cm(top), Cm(14.5), Cm(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = color
    accent.line.fill.background()

    tf = add_text(slide, left+1, top+0.4, 12.5, 5, title, size=17, color=color, bold=True)
    for item in items:
        add_para(tf, item, size=13, color=WHITE)
    add_para(tf, "", size=4)
    add_para(tf, f"  >> {diff}", size=12, color=GOLD, bold=True)

# ============================================================
# 슬라이드 12: 기술 스택
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "07  기술 스택")

stacks = [
    ("게임 엔진", "Unity 2D", "2D RPG 개발, 씬 관리, 애니메이션"),
    ("프로그래밍", "C#", "게임 로직, 매니저 패턴, 이벤트 시스템"),
    ("렌더링", "URP", "최적화된 2D 렌더링, 커스텀 셰이더"),
    ("백엔드", "Backend SDK", "유저 인증, 클라우드 저장, 채팅, 랭킹"),
    ("UI", "TextMesh Pro", "고품질 텍스트 렌더링"),
    ("연출", "Timeline", "가챠 연출, 인트로, 시네마틱"),
    ("데이터", "ScriptableObject + JSON", "게임 데이터 정의 및 세이브/로드"),
    ("최적화", "Object Pooling", "몬스터/총알/이펙트 재활용, GC 최소화"),
]

for idx, (category, tech, desc) in enumerate(stacks):
    y = 5.5 + idx * 1.5
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(2), Cm(y), Cm(5.5), Cm(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x00, 0x7A, 0xCC)
    box.line.fill.background()
    tf_box = box.text_frame
    tf_box.paragraphs[0].text = category
    tf_box.paragraphs[0].font.size = Pt(12)
    tf_box.paragraphs[0].font.color.rgb = WHITE
    tf_box.paragraphs[0].font.bold = True
    tf_box.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text(slide, 8.5, y, 7, 1.2, tech, size=14, color=ACCENT_BLUE, bold=True)
    add_text(slide, 16, y, 17, 1.2, desc, size=12, color=LIGHT_GRAY)

# ============================================================
# 슬라이드 13: 핵심 기능 - 전투
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "08-1  핵심 기능 - 전투 시스템")

add_accent_bar(slide, 2, 5, 0.2, 12)
tf = add_text(slide, 3, 5, 15, 12, "실시간 2D 전투 시스템", size=20, color=GOLD, bold=True)
add_para(tf, "", size=6)
add_para(tf, "캐릭터 타입", size=16, color=ACCENT_BLUE, bold=True)
add_para(tf, "  근거리(Melee): 근접 공격 + 쿨다운", size=14, color=WHITE)
add_para(tf, "  원거리(Ranged): 자동 조준 발사", size=14, color=WHITE)
add_para(tf, "", size=6)
add_para(tf, "크리티컬 시스템", size=16, color=ACCENT_BLUE, bold=True)
add_para(tf, "  4단계: Normal / Rare / Epic / Legendary", size=14, color=WHITE)
add_para(tf, "  단계별 데미지 배율 & 시각 효과", size=14, color=WHITE)
add_para(tf, "", size=6)
add_para(tf, "몬스터 AI & 스폰", size=16, color=ACCENT_BLUE, bold=True)
add_para(tf, "  플레이어 감지 -> 추적 -> 공격 패턴", size=14, color=WHITE)
add_para(tf, "  WaveSpawner 웨이브 기반 스폰", size=14, color=WHITE)
add_para(tf, "", size=6)
add_para(tf, "오브젝트 풀링", size=16, color=ACCENT_BLUE, bold=True)
add_para(tf, "  PoolManager로 GC 최소화", size=14, color=WHITE)

# Right - flow
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(20), Cm(5), Cm(12), Cm(12))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3E)
box.line.color.rgb = RGBColor(0x44, 0x44, 0x66)

tf = add_text(slide, 21, 5.5, 10, 11, "전투 흐름", size=18, color=ACCENT_BLUE, bold=True)
add_para(tf, "", size=6)
flow_items = [
    "1. 이동 + 자동 조준",
    "2. 자동/수동 공격",
    "3. 크리티컬 판정 & 데미지",
    "4. DamagePopup 표시",
    "5. 몬스터 처치",
    "6. 골드/경험치/아이템 획득",
    "7. GameManager 이벤트 전파",
]
for item in flow_items:
    add_para(tf, f"  {item}", size=14, color=WHITE)
    add_para(tf, "     |", size=10, color=ACCENT_BLUE)

# ============================================================
# 슬라이드 14: 핵심 기능 - 무한의탑 & 월드보스레이드
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "08-2  핵심 기능 - 무한의탑 & 월드보스 레이드")

# 무한의탑 - Left
add_accent_bar(slide, 2, 5, 0.2, 6.5)
tf = add_text(slide, 3, 5, 15, 6.5, "무한의탑", size=22, color=GOLD, bold=True)
add_para(tf, "", size=6)
add_para(tf, "컨셉", size=16, color=ACCENT_BLUE, bold=True)
add_para(tf, "  끝없이 올라가는 층수 도전형 엔드 콘텐츠", size=14, color=WHITE)
add_para(tf, "  층수가 올라갈수록 몬스터 강화 & 보상 증가", size=14, color=WHITE)
add_para(tf, "", size=6)
add_para(tf, "시스템", size=16, color=ACCENT_BLUE, bold=True)
add_para(tf, "  매 층마다 강력한 몬스터 웨이브 등장", size=14, color=WHITE)
add_para(tf, "  특정 층 클리어 시 특별 보상 (장비/보석)", size=14, color=WHITE)
add_para(tf, "  최고 기록 랭킹으로 경쟁 요소 제공", size=14, color=WHITE)
add_para(tf, "  장비/스킬 성장의 목표 콘텐츠", size=14, color=WHITE)

# 월드보스레이드 - Left bottom
add_accent_bar(slide, 2, 12.5, 0.2, 5)
tf = add_text(slide, 3, 12.5, 15, 5, "월드보스 레이드", size=22, color=GOLD, bold=True)
add_para(tf, "", size=6)
add_para(tf, "컨셉", size=16, color=ACCENT_BLUE, bold=True)
add_para(tf, "  전체 서버 유저가 함께 도전하는 대형 보스", size=14, color=WHITE)
add_para(tf, "", size=4)
add_para(tf, "시스템", size=16, color=ACCENT_BLUE, bold=True)
add_para(tf, "  제한 시간 내 최대 데미지 딜링 경쟁", size=14, color=WHITE)
add_para(tf, "  기여도 기반 보상 분배", size=14, color=WHITE)
add_para(tf, "  주기적 보스 교체 (다양한 패턴)", size=14, color=WHITE)

# Right - 콘텐츠 사이클
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(20), Cm(5), Cm(12), Cm(12.5))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3E)
box.line.color.rgb = RGBColor(0x44, 0x44, 0x66)

tf = add_text(slide, 21, 5.5, 10, 12, "엔드 콘텐츠 사이클", size=18, color=ACCENT_BLUE, bold=True)
add_para(tf, "", size=8)
add_para(tf, "  장비 수집 (가챠)", size=15, color=WHITE)
add_para(tf, "       |", size=12, color=ACCENT_BLUE)
add_para(tf, "  장비 강화 (30레벨)", size=15, color=WHITE)
add_para(tf, "       |", size=12, color=ACCENT_BLUE)
add_para(tf, "  무한의탑 도전", size=15, color=GOLD, bold=True)
add_para(tf, "       |", size=12, color=ACCENT_BLUE)
add_para(tf, "  더 높은 층 = 더 좋은 보상", size=15, color=WHITE)
add_para(tf, "       |", size=12, color=ACCENT_BLUE)
add_para(tf, "  월드보스 레이드 참여", size=15, color=GOLD, bold=True)
add_para(tf, "       |", size=12, color=ACCENT_BLUE)
add_para(tf, "  기여도 보상 획득", size=15, color=WHITE)
add_para(tf, "       |", size=12, color=ACCENT_BLUE)
add_para(tf, "  더 강한 장비로 반복!", size=15, color=GREEN, bold=True)

# ============================================================
# 슬라이드 15: 핵심 기능 - 가챠 & 장비
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "08-3  핵심 기능 - 가챠 & 장비 강화")

add_accent_bar(slide, 2, 5, 0.2, 5.5)
tf = add_text(slide, 3, 5, 14, 5.5, "가챠 시스템", size=20, color=GOLD, bold=True)
add_para(tf, "", size=4)
add_para(tf, "  장비 가챠: 5티어 확률 풀", size=14, color=WHITE)
add_para(tf, "  동료 가챠: 멀티탭 소환", size=14, color=WHITE)
add_para(tf, "  재화: 티켓 + 보석(Gem)", size=14, color=WHITE)
add_para(tf, "  Timeline 연출 (레전더리 특수)", size=14, color=WHITE)
add_para(tf, "  중복 아이템 그룹핑 표시", size=14, color=WHITE)

add_accent_bar(slide, 2, 11.5, 0.2, 5.5)
tf = add_text(slide, 3, 11.5, 14, 5.5, "장비 강화 시스템", size=20, color=GOLD, bold=True)
add_para(tf, "", size=4)
add_para(tf, "  6종: 무기(좌/우), 투구, 갑옷, 장갑, 부츠", size=14, color=WHITE)
add_para(tf, "  스탯: 공격력, 방어력, HP, 속도, 크리율", size=14, color=WHITE)
add_para(tf, "  강화: 최대 30레벨 (성공률 감소)", size=14, color=WHITE)
add_para(tf, "  실패 시 레벨 하락/파괴 위험", size=14, color=WHITE)
add_para(tf, "  안전 강화 옵션 제공", size=14, color=WHITE)

# Right - rates
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(19), Cm(5), Cm(13), Cm(12))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3E)
box.line.color.rgb = RGBColor(0x44, 0x44, 0x66)

tf = add_text(slide, 20, 5.5, 11, 11, "강화 성공률", size=16, color=ACCENT_BLUE, bold=True)
add_para(tf, "", size=4)
rates = [
    ("Lv 1~5", "100%", GREEN),
    ("Lv 6~10", "80%", GREEN),
    ("Lv 11~15", "60%", GOLD),
    ("Lv 16~20", "45%", GOLD),
    ("Lv 21~25", "30%", RED),
    ("Lv 26~30", "20%", RED),
]
for lvl, rate, color in rates:
    add_para(tf, f"  {lvl}  ->  {rate}", size=14, color=color)
    add_para(tf, "", size=4)

add_para(tf, "", size=6)
add_para(tf, "강화 비용", size=16, color=ACCENT_BLUE, bold=True)
add_para(tf, "  골드 + 농장 포인트(CropPoints)", size=14, color=WHITE)

# ============================================================
# 슬라이드 15: 핵심 기능 - 농장 & 스킬
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "08-4  핵심 기능 - 농장 & 스킬")

add_accent_bar(slide, 2, 5, 0.2, 6)
tf = add_text(slide, 3, 5, 14, 6, "농장 시스템 (레알팜 참고)", size=20, color=GOLD, bold=True)
add_para(tf, "", size=4)
add_para(tf, "  작물 심기/재배/수확 (성장 단계)", size=14, color=WHITE)
add_para(tf, "  건물 건설 & 업그레이드", size=14, color=WHITE)
add_para(tf, "  비료/물 시스템 (성장 부스트)", size=14, color=WHITE)
add_para(tf, "  경매장: 다른 플레이어와 거래", size=14, color=WHITE)
add_para(tf, "  농장 전용 재화 (CropPoints)", size=14, color=WHITE)

add_accent_bar(slide, 2, 12, 0.2, 5)
tf = add_text(slide, 3, 12, 14, 5, "스킬 시스템", size=20, color=GOLD, bold=True)
add_para(tf, "", size=4)
add_para(tf, "  포인트 기반 스킬 트리 학습", size=14, color=WHITE)
add_para(tf, "  핫바 드래그 앤 드롭 배치 (8슬롯)", size=14, color=WHITE)
add_para(tf, "  장비 패시브 스킬 자동 적용", size=14, color=WHITE)
add_para(tf, "  쿨다운 관리 + 자동 캐스트 토글", size=14, color=WHITE)

# Right - scene
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(19), Cm(5), Cm(13), Cm(12))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3E)
box.line.color.rgb = RGBColor(0x44, 0x44, 0x66)

tf = add_text(slide, 20, 5.5, 11, 11, "씬 구조", size=18, color=ACCENT_BLUE, bold=True)
add_para(tf, "", size=6)
scenes = [
    ("IntroVideo", "스플래시 / 인트로"),
    ("LoginScene", "로그인 + 캐릭터 슬롯"),
    ("MainScene", "전투, 인벤, 가챠, 상점"),
    ("FarmScene", "농장, 건물, 경매"),
]
for name, desc in scenes:
    add_para(tf, f"  {name}", size=15, color=ACCENT_BLUE, bold=True)
    add_para(tf, f"    -> {desc}", size=13, color=WHITE)
    add_para(tf, "", size=6)

add_para(tf, "데이터 흐름", size=18, color=ACCENT_BLUE, bold=True)
add_para(tf, "", size=4)
add_para(tf, "  GameDataBridge (Static)", size=13, color=WHITE)
add_para(tf, "  -> JSON 직렬화 (암호화)", size=13, color=WHITE)
add_para(tf, "  -> Backend 클라우드 백업", size=13, color=WHITE)

# ============================================================
# 슬라이드 16: 치매예방 미니게임
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "09  치매예방 미니게임", "두뇌 훈련 콘텐츠")

add_text(slide, 2, 4.2, 30, 1.2, "게임 내 미니게임 형태로 치매예방 두뇌 훈련 콘텐츠를 제공합니다", size=15, color=LIGHT_GRAY)

mini_games = [
    ("기억력 카드 뒤집기", [
        "  카드를 뒤집어 같은 그림 찾기",
        "  난이도별 카드 수 증가",
        "  제한 시간 내 모두 맞추기",
        "  기억력 & 집중력 향상",
    ], ACCENT_BLUE),
    ("숫자 순서 맞추기", [
        "  화면에 표시된 숫자를 순서대로 터치",
        "  단계별 숫자 개수 증가",
        "  역순 모드 (큰 수 -> 작은 수)",
        "  작업 기억력 훈련",
    ], GOLD),
    ("색상 판별 게임", [
        "  글자의 '색상'을 빠르게 판별",
        "  스트룹 효과 기반 두뇌 훈련",
        "  반응 속도 & 인지 유연성 향상",
        "  점수 기반 랭킹 경쟁",
    ], GREEN),
    ("패턴 기억하기", [
        "  화면에 나타난 패턴을 기억 후 재현",
        "  단계별 패턴 복잡도 증가",
        "  시각적 기억력 강화",
        "  전연령 즐길 수 있는 난이도",
    ], ACCENT_PURPLE),
]

for idx, (title, items, color) in enumerate(mini_games):
    col = idx % 2
    row = idx // 2
    left = 2 + col * 16.5
    top = 5.5 + row * 6

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(left), Cm(top), Cm(14.5), Cm(5.3))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3E)
    box.line.color.rgb = RGBColor(0x44, 0x44, 0x66)

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(left), Cm(top), Cm(14.5), Cm(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = color
    accent.line.fill.background()

    tf = add_text(slide, left+1, top+0.4, 12.5, 4.5, title, size=18, color=color, bold=True)
    for item in items:
        add_para(tf, item, size=13, color=WHITE)

# ============================================================
# 슬라이드 17: 아키텍처 & 설계 패턴
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "10  아키텍처 & 설계 패턴")

patterns_data = [
    ("싱글톤 매니저", "씬-로컬 싱글톤 패턴\nInstance 중복 방지", ACCENT_BLUE),
    ("데이터 브릿지", "GameDataBridge (Static)\n씬 간 데이터 전달 + JSON 직렬화", GOLD),
    ("이벤트 시스템", "OnGoldChanged, OnExpChanged 등\n디커플링된 UI 업데이트", GREEN),
    ("오브젝트 풀링", "PoolManager\n몬스터/총알/이펙트 재활용", ORANGE),
    ("실행 순서 제어", "GameManager(-100)\nSaveLoadManager(+100)", ACCENT_PURPLE),
    ("Backend 연동", "BackendManager\n인증, 저장, 채팅, 랭킹 통합", RED),
]

for idx, (name, desc, color) in enumerate(patterns_data):
    col = idx % 3
    row = idx // 3
    left = 2 + col * 11
    top = 5 + row * 6.5

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(left), Cm(top), Cm(9.5), Cm(5.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3E)
    box.line.color.rgb = RGBColor(0x44, 0x44, 0x66)

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(left), Cm(top), Cm(9.5), Cm(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = color
    accent.line.fill.background()

    tf = add_text(slide, left+0.8, top+0.5, 8, 4.5, name, size=18, color=color, bold=True)
    for line in desc.split("\n"):
        add_para(tf, line, size=12, color=WHITE)

# ============================================================
# 슬라이드 17: QA 결과
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "11-1  QA 결과와 수정 사항")

bugs = [
    ("스크롤 튀는 현상",
     "UI 스크롤뷰 터치 시 위치가 갑자기 튀는 버그",
     "ScrollRect content 위치 보정 로직 추가"),
    ("가챠-동료 패널 연동 해제",
     "가챠/동료 패널이 잘못 연결, 한쪽 닫으면 다른 쪽도 해제",
     "패널 간 독립적 참조 관리로 수정"),
    ("이미지 매칭 오류",
     "장비/아이템 이미지가 올바르게 매칭되지 않음",
     "정확한 키 매핑 로직으로 변경"),
    ("씬 전환 시 전투 비활성화",
     "씬 전환 중 전투 계속 진행으로 에러 발생",
     "씬 전환 시 전투 상태 명시적 비활성화"),
]

for idx, (title, desc, fix) in enumerate(bugs):
    y = 5 + idx * 3.2
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(2), Cm(y), Cm(30), Cm(2.8))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3E)
    box.line.color.rgb = RGBColor(0x44, 0x44, 0x66)

    tf = add_text(slide, 3, y+0.2, 28, 2.3, f"[BUG] {title}", size=16, color=RED, bold=True)
    add_para(tf, f"  {desc}", size=13, color=LIGHT_GRAY)
    add_para(tf, f"  -> {fix}", size=13, color=GREEN)

# ============================================================
# 슬라이드 18: QA 개선
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "11-2  QA 결과와 개선 사항")

improvements = [
    ("튜토리얼 중 버튼 차단",
     "튜토리얼 중 포커스 외 버튼 클릭 시 의도치 않은 동작",
     "ShouldBlockNonFocusButtons로 포커스 외 버튼 차단"),
    ("UI 클릭 가드",
     "빠른 연속 클릭으로 중복 동작 발생",
     "UIClickGuard 전역 클릭 쿨다운 패턴 적용"),
    ("가챠 결과 그룹핑",
     "같은 아이템이 개별 표시되어 가독성 저하",
     "중복 아이템 그룹핑 패턴으로 UI 개선"),
]

for idx, (title, feedback, improvement) in enumerate(improvements):
    y = 5 + idx * 4
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(2), Cm(y), Cm(30), Cm(3.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3E)
    box.line.color.rgb = RGBColor(0x44, 0x44, 0x66)

    tf = add_text(slide, 3, y+0.2, 28, 3, f"[IMPROVE] {title}", size=16, color=ACCENT_BLUE, bold=True)
    add_para(tf, f"  {feedback}", size=13, color=LIGHT_GRAY)
    add_para(tf, f"  -> {improvement}", size=13, color=GREEN)

# ============================================================
# 슬라이드 19: 시연 동영상
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "12  시연 동영상")

box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(4), Cm(5), Cm(26), Cm(11))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3E)
box.line.color.rgb = RGBColor(0x44, 0x44, 0x66)

tf = add_text(slide, 6, 7, 22, 8, "시연 동영상", size=32, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
add_para(tf, "", size=12)
add_para(tf, "별도 파일로 제출", size=20, color=WHITE, align=PP_ALIGN.CENTER)
add_para(tf, "", size=8)
add_para(tf, "용량: 100MB 이하  |  길이: 5~10분", size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_para(tf, "기능별 소개 음성 포함", size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 슬라이드 20: 감사합니다
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

accent_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Cm(1), prs.slide_height)
accent_rect.fill.solid()
accent_rect.fill.fore_color.rgb = ACCENT_BLUE
accent_rect.line.fill.background()

add_text(slide, 4, 5, 26, 4, "감사합니다", size=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
tf = add_text(slide, 4, 10, 26, 3, "INFINITY LEVEL UP", size=28, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
add_para(tf, "K-Digital Training 팀별 프로젝트 결과보고서", size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Cm(18.5), prs.slide_width, Cm(0.6))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_BLUE
bar.line.fill.background()

# ============================================================
# SAVE
# ============================================================
output_path = 'INFINITY_LEVELUP_결과보고서.pptx'
prs.save(output_path)
print(f'PPT 생성 완료: {output_path}')
print(f'총 슬라이드: {len(prs.slides)}장')
